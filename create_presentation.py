from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

def apply_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(21, 40, 40)  # Dark teal background

def add_logo(slide):
    # Add Codeium logo placeholder in top-left corner
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(1)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(124, 237, 237)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = "codeium"
    tf.paragraphs[0].font.color.rgb = RGBColor(21, 40, 40)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Using title slide layout
    apply_slide_background(slide)
    add_logo(slide)
    
    # Center the title and subtitle more precisely
    title_shape = slide.shapes.title
    title_shape.top = Inches(2.5)
    title_shape.left = Inches(1)
    title_shape.width = Inches(8)
    title_shape.text = title
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.top = Inches(4)
    subtitle_shape.left = Inches(1)
    subtitle_shape.width = Inches(8)
    subtitle_shape.text = subtitle
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(124, 237, 237)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Style the subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(124, 237, 237)
    subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_content_slide(prs, title, content, layout_index=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])  # Using content layout
    apply_slide_background(slide)
    add_logo(slide)
    
    # Add and style title with better positioning
    title_shape = slide.shapes.title
    title_shape.top = Inches(0.8)
    title_shape.left = Inches(1)
    title_shape.width = Inches(8)
    title_shape.text = title
    
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(124, 237, 237)
    
    # Add content with better positioning
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(124, 237, 237)
        p.level = 0 if not item.startswith('  ') else 1

def create_presentation():
    prs = Presentation()
    
    # Set default slide size to 16:9 (widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, 
                   "Windsurf Cascade + UKG",
                   "Transforming Developer Productivity with AI-Powered Code Assistance")
    
    # Enterprise Success
    enterprise_content = [
        "JP Morgan Chase Success:",
        "  • 10,000+ Active Developers",
        "  • 45% Productivity Increase",
        "",
        "Dell Technologies Impact:",
        "  • 3x Faster Code Reviews",
        "  • Significant improvement in code quality",
        "  • Reduced cognitive load on development teams"
    ]
    add_content_slide(prs, "Enterprise Success Stories", enterprise_content)
    
    # Industry Recognition
    recognition_content = [
        "Gartner® Magic Quadrant™ 2024:",
        "  • Leader in AI-Assisted Software Engineering",
        "",
        "Forrester Wave™ 2024:",
        "  • Strong Performer in Developer Experience Platforms",
        "",
        "IDC MarketScape 2024:",
        "  • Leader in AI-Powered Development Tools"
    ]
    add_content_slide(prs, "Industry Recognition", recognition_content)
    
    # ROI Analysis
    roi_content = [
        "Large-Scale Migration:",
        "  • Traditional: 10 hours ($1000)",
        "  • With Cascade: 2 hours ($200)",
        "  • 5x Cost Reduction",
        "",
        "Code Understanding:",
        "  • Traditional: 2 hours ($200)",
        "  • With Cascade: 6 minutes ($10)",
        "  • 20x Cost Reduction"
    ]
    add_content_slide(prs, "ROI Analysis", roi_content)
    
    # Key Benefits
    benefits_content = [
        "Enhanced Productivity:",
        "  • Faster code completion",
        "  • Reduced context switching",
        "  • Automated documentation",
        "",
        "Code Quality:",
        "  • Consistent coding standards",
        "  • Early bug detection",
        "  • Best practice suggestions"
    ]
    add_content_slide(prs, "Key Benefits", benefits_content)
    
    # Success Metrics
    metrics_content = [
        "Key Performance Indicators:",
        "  • 40% Reduction in Development Time",
        "  • 50% Fewer Code Review Iterations",
        "  • 30% Decrease in Bug Reports",
        "",
        "Developer Experience:",
        "  • Improved code quality",
        "  • Faster onboarding",
        "  • Reduced cognitive load"
    ]
    add_content_slide(prs, "Success Metrics", metrics_content)
    
    # Save the presentation
    prs.save('windsurf_cascade_ukg.pptx')

if __name__ == "__main__":
    create_presentation()
